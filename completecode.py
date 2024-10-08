import os
import asyncio
from openai import OpenAI

client = OpenAI(api_key='')
from crawl4ai import AsyncWebCrawler
from pptx import Presentation

# Step 1: Set your OpenAI API Key
  # You can replace this with your actual API key

# Asynchronous function to fetch data using Crawl4AI and save it to a file
async def fetch_data_and_save_to_file(url, filename="crawled_content.txt"):
    async with AsyncWebCrawler(verbose=True) as crawler:
        result = await crawler.arun(url=url)

        # Extract the full content from the result
        full_content = result.markdown  # Assuming the content is in markdown format

        # Save the content to a file
        with open(filename, 'w', encoding='utf-8') as f:
            f.write(full_content)
        print(f"Content saved to {filename}")

# Wrapper to run asynchronous crawling and saving
def save_crawled_content(url, filename="crawled_content.txt"):
    asyncio.run(fetch_data_and_save_to_file(url, filename))

# Step 2: Function to summarize content using GPT-4 model
def summarize_content_gpt4(content, topic):
    prompt = f"Extract and summarize the content related to {topic} from the following text:\n\n{content}\n\nPlease provide a concise and professional summary."

    response = client.chat.completions.create(model="gpt-4",  # GPT-4 model for ChatCompletion API
    messages=[
        {"role": "system", "content": "You are an expert assistant for summarization."},
        {"role": "user", "content": prompt}
    ],
    max_tokens=500,
    temperature=0.5)

    summary = response.choices[0].message.content.strip()
    return summary

# Step 3: Function to extract and summarize content for each topic from the file
def extract_and_summarize_from_file(filename, topic):
    with open(filename, 'r', encoding='utf-8') as f:
        content = f.read()

    # Call GPT-4 API to summarize the topic content
    summarized_content = summarize_content_gpt4(content, topic)
    return summarized_content

# Step 4: Function to create a PowerPoint presentation based on topics
def create_presentation(template_path, topics_content):
    prs = Presentation(template_path)  # Load the provided template

    for topic, content in topics_content.items():
        slide_layout = prs.slide_layouts[1]  # Layout 1: 'Title and Content'
        slide = prs.slides.add_slide(slide_layout)

        # Set slide title
        title = slide.shapes.title
        title.text = topic

        # Set slide content
        content_placeholder = slide.shapes.placeholders[1]
        content_placeholder.text = content

    prs.save('generated_presentation.pptx')
    print("Presentation created: 'generated_presentation.pptx'")

# Step 5: Main function to integrate all steps
def generate_presentation():
    # URL to crawl content from
    url = "https://en.wikipedia.org/wiki/Artificial_intelligence_in_healthcare"

    # Step 1: Crawl and save content to a file
    filename = "crawled_content.txt"
    save_crawled_content(url, filename)

    # Step 2: Define the topics and extract summarized content for each one
    topics = ['Dermatology', 'Gastroenterology', 'Neurology']
    topics_content = {topic: extract_and_summarize_from_file(filename, topic) for topic in topics}

    # Step 3: Create the PowerPoint presentation using the summarized content
    create_presentation("my_template.pptx", topics_content)

# Step 6: Run the process to generate the presentation
generate_presentation()
