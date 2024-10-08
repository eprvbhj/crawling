import asyncio
from crawl4ai import AsyncWebCrawler
from pptx import Presentation
from transformers import pipeline

# Initialize the Hugging Face summarization pipeline
summarizer = pipeline("summarization")

# Asynchronous function to fetch data using Crawl4AI
async def fetch_data_with_crawl4ai(url):
    async with AsyncWebCrawler(verbose=True) as crawler:
        result = await crawler.arun(url=url)
        return result.markdown

# Function to summarize content using a Transformer model (from Hugging Face)
def summarize_content(content, max_length=500):
    # Limit content length to 1024 tokens before sending to the summarizer
    if len(content) > 1024:
        content = content[:1024]
    
    # Use the transformer to summarize if the content is too long
    summary = summarizer(content, max_length=max_length, min_length=100, do_sample=False)
    return summary[0]['summary_text']

# Function to extract relevant content for each topic
def extract_topic_content(content, topic):
    # Simple keyword-based search for content related to the topic
    topic_index = content.lower().find(topic.lower())
    if topic_index == -1:
        return f"No content found for {topic}"

    # Extract content starting from the found topic
    extracted_content = content[topic_index:topic_index + 2000]  # Limit content for each topic
    return summarize_content(extracted_content)

# Synchronous wrapper to call asynchronous code
def get_crawled_data(url):
    raw_content = asyncio.run(fetch_data_with_crawl4ai(url))
    return raw_content

# Function to create a PowerPoint presentation based on topics
def create_presentation(template_path, topics_content):
    prs = Presentation(template_path)  # Load the provided template

    for topic, content in topics_content.items():
        slide_layout = prs.slide_layouts[1]  # Layout 1: 'Title and Content'
        slide = prs.slides.add_slide(slide_layout)

        # Set slide title
        title = slide.shapes.title
        title.text = topic

        # Set slide content
        content_placeholder = slide.shapes.placeholders[1]
        content_placeholder.text = content

    prs.save('generated_presentation.pptx')
    print("Presentation created: 'generated_presentation.pptx'")

# Main function to generate the presentation
def generate_presentation():
    url = "https://en.wikipedia.org/wiki/Artificial_intelligence_in_healthcare"

    # Fetch and summarize content from the URL
    raw_content = get_crawled_data(url)

    # Define the topics and extract content for each one
    topics = ['Dermatology', 'Gastroenterology', 'Neurology']
    topics_content = {topic: extract_topic_content(raw_content, topic) for topic in topics}

    # Create the PowerPoint presentation with the provided template
    create_presentation("my_template.pptx", topics_content)

# Run the presentation generation with the given topics
generate_presentation()
